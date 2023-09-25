import streamlit as st
import boto3
import PyPDF2
import langchain
from odf import load as odf_load
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS

def read_data(s3_bucket, s3_key):
    s3 = boto3.client('s3')
    response = s3.get_object(Bucket=s3_bucket, Key=s3_key)

    file_extension = s3_key.split('.')[-1]

    if file_extension == 'pdf':
        pdf_reader = PyPDF2.PdfFileReader(response['Body'])
        text = ''
        for page in range(pdf_reader.getNumPages()):
            text += pdf_reader.getPage(page).extractText()
    elif file_extension == 'odp' or file_extension == 'odt':
        load = odf_load(response['Body'])
        text = load.text
    elif file_extension == 'docx':
        document = Document(response['Body'])
        text = ''
        for paragraph in document.paragraphs:
            text += paragraph.text
    elif file_extension == 'pptx':
        presentation = Presentation(response['Body'])
        text = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
                            text += ' '
    elif file_extension == 'txt':
        text = response['Body'].read().decode()
    else:
        text = ''

    return text

def append_data_to_pdf(pdf_writer, text):
    pdf_writer.addPage(PyPDF2.PdfTextFormField(text))

quiz_data = [
    {
        "question": "What is the capital of France?",
        "options": ["Paris", "Berlin", "Madrid", "London"],
        "correct_option": 0,
    },
    {
        "question": "Which planet is known as the Red Planet?",
        "options": ["Venus", "Mars", "Jupiter", "Saturn"],
        "correct_option": 1,
    }
]

def main():
    st.title("📝 MITZ Session Feedback AI")

    # pdf_writer = PyPDF2.PdfFileWriter()
    
    # s3 = boto3.client('s3')
    # objects = s3.list_objects(Bucket='my-bucket', Prefix='my-folder')['Contents']
    # for obj in objects:
    #     s3_key = obj['Key']
    #     text = read_data('my-bucket', s3_key)
    #     append_data_to_pdf(pdf_writer, text)

    # pdf_object_string = pdf_writer.getData()

    # with open('merged_pdf.pdf', 'wb') as f:
    #     f.write(pdf_object_string)

    ###written everything to one pdf file

#read everything and create embedding inside the vectorDB

# def get_pdf_text(pdf_docs):
#     text = ""
#     for pdf in pdf_docs:
#         pdf_reader = PdfReader(pdf)
#         for page in pdf_reader.pages:
#             text += page.extract_text()
#         return text
        

# def get_text_chunks(text):
#     text_splitter = CharacterTextSplitter(
#         separator="\n",
#         chunk_size = 1000,
#         chunk_overlap = 200, #to contain full sentences to not miss out the meaning
#         length_function = len
#     )
#     chunks = text_splitter.split_text(text)
#     return chunks


# def get_vectorstore(text_chunks):
#     embeddings = OpenAIEmbeddings() ## we can use other embedding models but its way slower
#     vectorstore = FAISS.from_texts(texts=text_chunks,embedding=embeddings)
#     return vectorstore

# vectorstore = get_vectorstore(text_chunks)

# retrieval_chain = langchain.Retrieval(vector_store=vectorstore)

# prompt = f""Create 10 mcqs in the format{quiz_data}""

# response = retrieval_chain.run(prompt) /// we can use response as the quizdata list of dict for our application

    question_flags = [False] * len(quiz_data)

    selected_options = [-1] * len(quiz_data)

    for i, question_data in enumerate(quiz_data, start=1):
        st.subheader(f"Question {i}: {question_data['question']}")

        selectbox_key = f"q{i}_selectbox"

        selected_option = st.selectbox("Choose an option:", [""] + question_data['options'], key=selectbox_key)

        if selected_option:
            selected_option_index = question_data['options'].index(selected_option)
            selected_options[i - 1] = selected_option_index

    if st.button("Submit"):
        for i, question_data in enumerate(quiz_data):
            if selected_options[i] == question_data['correct_option']:
                question_flags[i] = True

        st.write(f"Your quiz score is: {question_flags.count(True)}")

if __name__ == "__main__":
    main()
